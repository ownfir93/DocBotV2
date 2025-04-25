#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_index_pipeline.py – Creates a vector store index from source documents in GCS.
"""

# Standard imports
import os
import sys
import time
from pathlib import Path
import tempfile
import json
import atexit
import logging

# --- START: Replit Service Account Key Handling ---
# Check if running in Replit and the JSON content secret is set
gcp_json_key_content = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS_JSON")
temp_key_file_path = None

if gcp_json_key_content:
    try:
        # Create a temporary file to store the key
        key_file_descriptor, temp_key_file_path = tempfile.mkstemp(suffix=".json")
        print(f"[INFO] Creating temporary key file at: {temp_key_file_path}")

        with os.fdopen(key_file_descriptor, 'w') as temp_key_file:
            temp_key_file.write(gcp_json_key_content)

        # Set the environment variable Google libraries expect
        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = temp_key_file_path
        print(f"[INFO] GOOGLE_APPLICATION_CREDENTIALS set to: {temp_key_file_path}")

        # Register a cleanup function
        def cleanup_keyfile():
            if temp_key_file_path and os.path.exists(temp_key_file_path):
                print(f"[INFO] Cleaning up temporary key file: {temp_key_file_path}")
                try: os.remove(temp_key_file_path)
                except OSError as e: print(f"[WARN] Could not remove temp key file {temp_key_file_path}: {e}")
        atexit.register(cleanup_keyfile)

    except Exception as e:
        print(f"[ERROR] Failed to process GOOGLE_APPLICATION_CREDENTIALS_JSON secret: {e}", file=sys.stderr)
        # Exit if auth setup fails, as the rest of the script depends on it
        sys.exit(1)
else:
    # Check if GOOGLE_APPLICATION_CREDENTIALS is set directly (e.g., in Replit env)
    if not os.environ.get("GOOGLE_APPLICATION_CREDENTIALS"):
         print("[ERROR] Authentication configuration missing.", file=sys.stderr)
         print("Ensure GOOGLE_APPLICATION_CREDENTIALS_JSON secret is set in Replit.", file=sys.stderr)
         sys.exit(1)
    else:
         print("[INFO] Using GOOGLE_APPLICATION_CREDENTIALS environment variable directly.")


# --- END: Replit Service Account Key Handling ---

from dotenv import load_dotenv
# Load environment variables
load_dotenv()

# Continue with the rest of your imports and code...

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_index_pipeline.py - Standalone script to:
1. List source documents from a GCS prefix.
2. Download/read documents, convert to text.
3. Generate text embeddings using OpenAI.
4. Build a Chroma vector index locally, storing GCS path in metadata.
5. Upload the generated Chroma index directory back to GCS.

Reads sources FROM GCS, writes index TO GCS.
"""

import os
import sys
import json
import logging
import re
import shutil
import tempfile
from pathlib import Path
from typing import List, Dict, Optional, Tuple

from dotenv import load_dotenv

load_dotenv()

# Document Parsing
from pdfminer.high_level import extract_text
from pptx import Presentation
import docx2txt

# LangChain
try:
    from langchain_openai import OpenAIEmbeddings
    from langchain_chroma import Chroma
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_core.documents import Document
except ImportError as e:
    print(
        f"ERROR: LangChain packages not found: {e}. pip install -r requirements.txt"
    )
    sys.exit(1)

# Google Cloud
try:
    from google.cloud import storage
    from google.oauth2 import service_account
    from google.auth.exceptions import DefaultCredentialsError
except ImportError:
    print(
        "ERROR: google-cloud-storage not found. pip install google-cloud-storage"
    )
    sys.exit(1)
except DefaultCredentialsError:
    print(
        "ERROR: Google Cloud Default Credentials not found. Ensure `gcloud auth application-default login` or GOOGLE_APPLICATION_CREDENTIALS."
    )
    sys.exit(1)

# --- CONFIGURATION ---
# GCS Configuration (REQUIRED - Set in .env or environment)
GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME")
GCS_SOURCE_PREFIX = os.getenv(
    "GCS_SOURCE_PREFIX",
    "source_documents/")  # Folder IN BUCKET with source PDFs, DOCX etc.
GCS_INDEX_PATH = os.getenv(
    "GCS_INDEX_PATH",
    "vector_index/chroma_db/")  # Folder IN BUCKET to store the built index

# Ensure trailing slashes for prefixes if needed by logic later
GCS_SOURCE_PREFIX = GCS_SOURCE_PREFIX.strip(
    '/') + '/' if GCS_SOURCE_PREFIX else ""
GCS_INDEX_PATH = GCS_INDEX_PATH.strip('/') + '/' if GCS_INDEX_PATH else ""

# Local Temporary Paths (Pipeline runtime only)
TEMP_DOWNLOAD_DIR = Path(tempfile.mkdtemp(prefix="pipeline_downloads_"))
TEMP_CHROMA_DB_PATH = Path(tempfile.mkdtemp(prefix="pipeline_chroma_build_"))

# OpenAI & Embedding Config
# Ensure OPENAI_API_KEY is set in environment
EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL_NAME",
                                 "text-embedding-3-small")
TEXT_SPLIT_CHUNK_SIZE = int(os.getenv("TEXT_SPLIT_CHUNK_SIZE", 1000))
TEXT_SPLIT_CHUNK_OVERLAP = int(os.getenv("TEXT_SPLIT_CHUNK_OVERLAP", 200))

# --- LOGGING SETUP ---
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] [%(funcName)s:%(lineno)d] %(message)s")
# Quieten verbose libraries
for logger_name in [
        "openai", "httpx", "httpcore", "urllib3", "pdfminer", "chromadb",
        "google.cloud.storage", "google.auth", "googleapiclient"
]:
    logging.getLogger(logger_name).setLevel(logging.WARNING)
logging.getLogger("langchain").setLevel(logging.INFO)  # Keep Langchain info

# --- GLOBAL VARIABLES ---
embeddings_model = None
gcs_client = None
gcs_bucket = None


# --- GCS HELPER FUNCTIONS ---
def get_gcs_client_and_bucket():
    """Initializes GCS client and bucket object."""
    global gcs_client, gcs_bucket
    if gcs_client is None or gcs_bucket is None:
        if not GCS_BUCKET_NAME:
            logging.critical("GCS_BUCKET_NAME environment variable not set.")
            raise ValueError("GCS_BUCKET_NAME not configured.")
        try:
            key_file_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
            if not key_file_path or not os.path.exists(key_file_path):
                 raise ValueError("GOOGLE_APPLICATION_CREDENTIALS path not valid or file missing.")

            # Explicitly load credentials from the file
            credentials = service_account.Credentials.from_service_account_file(key_file_path)
            logging.info(f"Explicitly loaded credentials from: {key_file_path}")

            # Initialize client with explicit project AND credentials
            gcs_client = storage.Client(project="elevado", credentials=credentials)
            gcs_bucket = gcs_client.get_bucket(GCS_BUCKET_NAME)
            logging.info(f"GCS client initialized and bucket '{GCS_BUCKET_NAME}' connected.")
        except Exception as e:
            logging.critical(
                f"Failed to initialize GCS client or get bucket '{GCS_BUCKET_NAME}': {e}",
                exc_info=True)
            raise
    return gcs_client, gcs_bucket


def list_gcs_source_files(prefix: str) -> List[storage.Blob]:
    """Lists supported document files in a GCS prefix."""
    _client, bucket = get_gcs_client_and_bucket()
    supported_exts = {".pdf", ".pptx", ".docx", ".txt", ".md", ".html"}
    source_blobs = []
    logging.info(
        f"Scanning gs://{GCS_BUCKET_NAME}/{prefix} for source files...")

    try:
        blobs = list(
            bucket.list_blobs(prefix=prefix))  # List all blobs in prefix
        for blob in blobs:
            # Skip "directory" objects and check extension
            if not blob.name.endswith('/') and Path(
                    blob.name).suffix.lower() in supported_exts:
                source_blobs.append(blob)
        logging.info(
            f"Found {len(source_blobs)} potential source files in GCS.")
        return source_blobs
    except Exception as e:
        logging.error(
            f"Failed to list GCS blobs in gs://{GCS_BUCKET_NAME}/{prefix}: {e}",
            exc_info=True)
        raise


def download_blob_to_temp(blob: storage.Blob, temp_dir: Path) -> Path:
    """Downloads a GCS blob to a temporary local file."""
    local_path = temp_dir / Path(blob.name).name  # Use only filename locally
    try:
        logging.debug(
            f"Downloading gs://{blob.bucket.name}/{blob.name} to {local_path}")
        blob.download_to_filename(str(local_path))
        return local_path
    except Exception as e:
        logging.error(f"Failed to download blob {blob.name}: {e}",
                      exc_info=True)
        raise


def upload_directory_to_gcs(local_path: Path, gcs_destination_prefix: str):
    """Uploads the contents of a local directory to a GCS prefix."""
    _client, bucket = get_gcs_client_and_bucket()
    logging.info(
        f"Uploading directory '{local_path}' contents to GCS path 'gs://{GCS_BUCKET_NAME}/{gcs_destination_prefix}'..."
    )
    assert local_path.is_dir()
    files_uploaded_count = 0
    for local_file in local_path.rglob('*'):
        if local_file.is_file():
            relative_path = local_file.relative_to(local_path)
            blob_path = Path(gcs_destination_prefix.strip('/')).joinpath(
                relative_path).as_posix()
            logging.debug(
                f"Uploading {local_file} to gs://{GCS_BUCKET_NAME}/{blob_path}"
            )
            try:
                blob = bucket.blob(blob_path)
                blob.upload_from_filename(str(local_file))
                files_uploaded_count += 1
            except Exception as e:
                logging.error(
                    f"Failed to upload {local_file} to {blob_path}: {e}",
                    exc_info=True)
                # Decide if one failed upload should stop the whole process
                raise  # Option: Stop on first error
    logging.info(
        f"Successfully uploaded {files_uploaded_count} index files to GCS.")


# --- DOCUMENT PROCESSING FUNCTION ---
def process_source_blob(blob: storage.Blob,
                        temp_dir: Path) -> Optional[Document]:
    """Downloads blob, extracts text, creates Langchain Document with GCS path metadata."""
    local_file_path = None
    try:
        local_file_path = download_blob_to_temp(blob, temp_dir)
        text = ""
        file_lower = local_file_path.name.lower()

        # --- Extract Text (same logic as before, but using local_file_path) ---
        if file_lower.endswith(".pdf"):
            text = extract_text(str(local_file_path))
        elif file_lower.endswith(".pptx"):
            prs = Presentation(str(local_file_path))
            text = "\n\n".join(shape.text.strip() for slide in prs.slides
                               for shape in slide.shapes
                               if hasattr(shape, "text") and shape.text
                               and shape.text.strip())
        elif file_lower.endswith(".docx"):
            text = docx2txt.process(str(local_file_path))
        elif file_lower.endswith((".txt", ".md", ".html")):
            try:
                text = local_file_path.read_text("utf-8")
            except UnicodeDecodeError:
                text = local_file_path.read_text("latin-1", errors="replace")
        else:
            logging.warning(f"Skipping unsupported file type: {blob.name}")
            return None  # Skip this blob

        # --- Clean Text ---
        text = re.sub(r'\s{3,}', '\n\n', text).strip()
        if not text:
            logging.warning(f"No text extracted from {blob.name}. Skipping.")
            return None

        # --- Create Document with Metadata ---
        # Store the GCS path for linking later
        metadata = {
            "source_gcs_path":
            f"gs://{blob.bucket.name}/{blob.name}",
            "original_filename":
            Path(blob.name).name,
        }
        
        # --- START: Category Keyword Definitions ---
        # Based on GCS structure from gs://lwdocbot/Seismic Mirror/
        # Keys should match the 'normalized_type' generated in process_source_blob
        CATEGORY_KEYWORDS = {
            # Top Level or General
            "root": [], # Files directly under Seismic Mirror/
            "highspot_download": ["asset", "content", "downloaded file", "whitepaper", "report", "datasheet", "presentation", "competitive analysis", "proposal", "rfi response", "benchmark", "survey", "webinar"],
            "product_marketing": ["messaging", "positioning", "datasheet", "feature overview", "go-to-market", "competitive analysis", "presentation", "platform overview", "gartner", "demo", "analyst"],
            "one_pager": ["summary", "overview", "brief", "datasheet", "flyer", "handout", "sell sheet", "solution brief"],
            "deal_desk": ["pricing", "quoting", "approval", "contract", "SLA", "renewal", "SKU", "tiers", "purchasing", "legal"],
            "marketing_content": ["marketing asset", "campaign material"],
            "webinar": ["presentation", "recording", "online event", "demo", "marketing content"],

            # Customer Stories Structure
            "case_study": ["customer success", "success story", "client example", "implementation", "results", "ROI", "challenge", "solution", "b2b example", "use case", "client story", "testimonial"],
            # Note: Win reports also get case study keywords added below in the code logic if needed

            # Win Report Sub-Types (Can inherit from case_study if desired)
            "marketing_touchpoints_to_closed_won": ["marketing campaign", "lead generation", "attribution", "closed won", "customer journey", "deal analysis", "win report"],
            "who_what_win_report": ["deal summary", "competitive win", "sales success", "customer problem", "solution provided", "win report"],
            "win_report_ae": ["AE perspective", "sales success", "opportunity review", "client name", "competitive win", "win wire", "win report"],

            # RevOps & Enablement Structure
            "revops_enablement": ["revops", "sales enablement", "gtm operations", "process", "sales guidelines", "kpmg"], # Top level
            "enablement": ["training", "learning", "onboarding", "bootcamp", "glossary", "ICP", "TAL", "sales enablement", "product enablement", "faq"], # Sub-folder
            "meddicc": ["MEDDIC", "MEDDPICC", "sales process", "qualification", "customer engagement", "opportunity review"],
            "roi_calculator": ["ROI", "calculator", "value assessment", "business case", "maturity assessment", "cost savings", "value engineering"],
            "sales_deck": ["presentation", "pitch", "slides", "customer presentation", "NBM deck", "QBR", "overview", "sales collateral", "marketing deck"],

            # Value Framework Sub-Structure
            "discovery_guide": ["discovery questions", "qualification", "needs analysis", "pain points", "requirements gathering", "probing", "discovery process", "value framework", "worksheet", "template", "example"],
            "discovery_guide_example": ["discovery guide", "example", "template"], # Specific examples subfolder
            "discovery_guide_workshop": ["discovery guide", "workshop", "training", "worksheet", "framework"], # Specific workshops subfolder
            "value_landscape": ["value proposition", "market positioning", "value driver", "personas", "competitive landscape"],
            "value_map": ["value proposition", "value driver", "solution mapping", "benefits", "features"],

            # Fallback
            "unknown": [],
        }
        
        # --- END: Category Keyword Definitions ---

        # Determine document type based on path, prioritizing deeper paths
        # Remove the source prefix and potential leading/trailing slashes
        relative_path_str = blob.name
        if GCS_SOURCE_PREFIX and relative_path_str.startswith(GCS_SOURCE_PREFIX):
            relative_path_str = relative_path_str[len(GCS_SOURCE_PREFIX):]

        path_parts = Path(relative_path_str).parts
        doc_type = "unknown" # Default

        # Normalize function
        def normalize_key(text):
            text = text.lower()
            text = re.sub(r'[\(\)]', '', text) # Remove parentheses
            text = re.sub(r'[\+&]', 'and', text) # Replace + or & with 'and'
            text = re.sub(r'[^\w\s-]', '', text) # Remove non-alphanumeric except space, hyphen
            text = re.sub(r'\s+', '_', text).strip('_') # Replace spaces with underscore
            return text if text else "unknown"

        # Check specific known paths first (deepest first)
        full_path_lower = relative_path_str.lower()
        if "/case studies/win reports/marketing touchpoints to closed won/" in full_path_lower:
            doc_type = "marketing_touchpoints_to_closed_won"
        elif "/case studies/win reports/who, what, win reports/" in full_path_lower:
            doc_type = "who_what_win_report"
        elif "/case studies/win reports/win reports created by aes/" in full_path_lower:
            doc_type = "win_report_ae"
        elif "/case studies/" in full_path_lower:
             # Check if it's just a case study, not a sub-type of win report
             if "/win reports/" not in full_path_lower:
                  doc_type = "case_study"
             # If it IS under win reports but didn't match specific type, maybe assign general 'win_report'?
             # else:
             #     doc_type = "win_report" # Optional general win_report category
        elif "/value framework/discovery guide/discovery guide examples/" in full_path_lower:
            doc_type = "discovery_guide_example"
        elif "/value framework/discovery guide/discovery guide workshops/" in full_path_lower:
            doc_type = "discovery_guide_workshop"
        elif "/value framework/discovery guide/" in full_path_lower:
            doc_type = "discovery_guide"
        elif "/value framework/value landscape/" in full_path_lower:
            doc_type = "value_landscape"
        elif "/value framework/value map/" in full_path_lower:
            doc_type = "value_map"
        elif "/revops + enablement  (jaisy)/enablement/" in full_path_lower: # Note potential double space
            doc_type = "enablement"
        elif "/revops + enablement  (jaisy)/meddicc/" in full_path_lower:
            doc_type = "meddicc"
        elif "/revops + enablement  (jaisy)/roi calculators/" in full_path_lower:
            doc_type = "roi_calculator"
        elif "/revops + enablement  (jaisy)/sales decks/" in full_path_lower:
            doc_type = "sales_deck"
        elif "/revops + enablement  (jaisy)/" in full_path_lower: # Parent folder if not in subfolder
             doc_type = "revops_enablement"
        elif "/deal desk (natalie)/" in full_path_lower:
            doc_type = "deal_desk"
        elif "/highspot downloads/" in full_path_lower:
             doc_type = "highspot_download"
        elif "/marketing content (lila)/webinars/" in full_path_lower:
             doc_type = "webinar"
        elif "/marketing content (lila)/" in full_path_lower:
             doc_type = "marketing_content"
        elif "/one-pagers/" in full_path_lower:
             doc_type = "one_pager"
        elif "/product marketing/gartner demos 2025/" in full_path_lower:
             # Decide if demos get own type or inherit from product marketing
             doc_type = "product_marketing" # Example: Inherit
        elif "/product marketing/" in full_path_lower:
             doc_type = "product_marketing"
        elif len(path_parts) > 1 : # Use immediate parent if path has depth and no specific rule matched
            doc_type = normalize_key(path_parts[-2]) # Use parent folder name
        elif len(path_parts) == 1 and path_parts[0] != relative_path_str: # File in root?
             doc_type = "root"


        metadata['document_type'] = doc_type
        logging.debug(f"Assigned document_type: {doc_type} for blob {blob.name}")

        # Add category keywords based on the determined doc_type
        keywords = CATEGORY_KEYWORDS.get(doc_type, [])
        # --- USER REQUEST: Add case study keywords also to win reports ---
        if doc_type in ["win_report_ae", "who_what_win_report", "marketing_touchpoints_to_closed_won"]:
             keywords.extend(CATEGORY_KEYWORDS.get("case_study", []))
             keywords = list(set(keywords)) # Remove duplicates
        # --- END USER REQUEST ---
        if keywords:
            metadata['category_keywords'] = keywords
            logging.debug(f"Added keywords for type '{doc_type}': {keywords}")
        else:
            logging.debug(f"No specific keywords found for type '{doc_type}'.")

        logging.debug(f"Final metadata for blob {blob.name}: {metadata}")

        logging.debug(
            f"Processed blob {blob.name}, extracted {len(text)} chars. Assigned document_type: {doc_type}")
        return Document(page_content=text, metadata=metadata)

    except Exception as e:
        logging.error(f"Error processing blob {blob.name}: {e}", exc_info=True)
        return None  # Skip on error
    finally:
        # Clean up temporary downloaded file
        if local_file_path and local_file_path.exists():
            try:
                local_file_path.unlink()
            except OSError:
                logging.warning(
                    f"Could not delete temp file: {local_file_path}")


# --- INDEX BUILDING FUNCTION ---
def build_chroma_index(documents: List[Document],
                       output_db_path: Path) -> bool:
    """Builds Chroma DB locally from processed documents."""
    global embeddings_model
    if not embeddings_model:
        logging.critical("Embeddings model not initialized.")
        return False
    if not documents:
        logging.error("No documents provided to build index.")
        return False

    output_db_path_str = str(output_db_path)

    # Clean previous build attempt
    if output_db_path.exists():
        logging.warning(
            f"Removing existing temporary Chroma directory: {output_db_path_str}"
        )
        try:
            shutil.rmtree(output_db_path)
        except OSError as e:
            logging.error(f"Error removing outdated Chroma directory: {e}",
                          exc_info=True)
            return False

    try:
        # --- Split Documents ---
        logging.info(f"Splitting {len(documents)} documents...")
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=TEXT_SPLIT_CHUNK_SIZE,
            chunk_overlap=TEXT_SPLIT_CHUNK_OVERLAP,
            # Keep metadata when splitting
            keep_separator=False  # Often better for Chroma chunking
        )
        docs_split = text_splitter.split_documents(documents)
        logging.info(f"Split into {len(docs_split)} chunks.")
        if not docs_split:
            logging.error("No chunks generated after splitting.")
            return False

        # --- Create Chroma Database Locally ---
        logging.info(
            f"Creating temporary Chroma database at: {output_db_path_str}...")
        # This step performs the embedding and writing to local disk
        Chroma.from_documents(documents=docs_split,
                              embedding=embeddings_model,
                              persist_directory=output_db_path_str)
        logging.info(f"Temporary Chroma database created locally.")
        return True

    except Exception as e:
        logging.error(f"Failed during index building: {e}", exc_info=True)
        if output_db_path.exists():  # Clean up failed build
            try:
                shutil.rmtree(output_db_path)
            except OSError as rm_err:
                logging.error(f"Failed cleanup {output_db_path_str}: {rm_err}")
        return False


# --- MAIN PIPELINE EXECUTION ---
def main():
    """Runs the full GCS-based indexing pipeline."""
    logging.info("--- Starting GCS Data Indexing Pipeline ---")

    # --- Validate Configuration ---
    if not GCS_BUCKET_NAME or not GCS_INDEX_PATH or not GCS_SOURCE_PREFIX:
        logging.critical(
            "GCS_BUCKET_NAME, GCS_SOURCE_PREFIX, and GCS_INDEX_PATH must be set. Exiting."
        )
        sys.exit(1)
    if not os.getenv("OPENAI_API_KEY"):
        logging.critical("OPENAI_API_KEY not set. Exiting.")
        sys.exit(1)

    # --- Initialize Services ---
    global embeddings_model
    try:
        get_gcs_client_and_bucket()  # Init GCS first
        
        # Initialize Embeddings directly, letting it create its own client
        embeddings_model = OpenAIEmbeddings(model=EMBEDDING_MODEL_NAME)
        logging.info(f"OpenAI Embeddings model '{EMBEDDING_MODEL_NAME}' initialized.")

    except Exception as init_err:
        logging.critical(
            f"Service initialization failed: {init_err}. Exiting.",
            exc_info=True)
        sys.exit(1)

    all_docs: List[Document] = []
    processed_count = 0
    error_count = 0

    try:
        # --- Step 1: List Source Files from GCS ---
        source_blobs = list_gcs_source_files(GCS_SOURCE_PREFIX)
        if not source_blobs:
            logging.warning(
                "No source files found in GCS path. Pipeline finished.")
            return  # Exit cleanly

        # --- Step 2: Process Each Source File ---
        logging.info(f"Processing {len(source_blobs)} source blobs...")
        for blob in source_blobs:
            doc = process_source_blob(blob, TEMP_DOWNLOAD_DIR)
            if doc:
                all_docs.append(doc)
                processed_count += 1
            else:
                error_count += 1

        logging.info(
            f"Blob Processing Complete. Success: {processed_count}, Errors: {error_count}"
        )
        if not all_docs:
            logging.error(
                "No documents were successfully processed. Cannot build index."
            )
            return

        # --- Step 3: Build Chroma Index Locally ---
        if not build_chroma_index(all_docs, TEMP_CHROMA_DB_PATH):
            logging.error("Chroma index building failed. Aborting upload.")
            return

        # --- Step 4: Upload Index to GCS ---
        upload_directory_to_gcs(TEMP_CHROMA_DB_PATH, GCS_INDEX_PATH)

    except Exception as pipeline_err:
        logging.error(f"Pipeline execution failed: {pipeline_err}",
                      exc_info=True)
    finally:
        # --- Step 5: Cleanup Temporary Local Directories ---
        logging.info("Cleaning up temporary local directories...")
        try:
            shutil.rmtree(TEMP_DOWNLOAD_DIR)
            shutil.rmtree(TEMP_CHROMA_DB_PATH)
            logging.info("Temporary directories cleaned up.")
        except OSError as e:
            logging.warning(f"Could not remove temporary directories: {e}")

    logging.info("--- GCS Data Indexing Pipeline Finished ---")


if __name__ == "__main__":
    main()